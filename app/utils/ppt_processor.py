from pptx import Presentation
from pptx.util import Inches

def generate_ppt_from_template(template_path, slides_content, output_path):
    """
    Generates a new PPTX file using the uploaded template and slide contents.
    If template_path is None, creates a basic presentation.
    """
    if template_path:
        prs = Presentation(template_path)
        # Remove existing slides
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]
    else:
        # Create a basic presentation without template
        prs = Presentation()
        
    # Add new slides
    for slide_data in slides_content:
        if len(prs.slide_layouts) > 1:
            layout = prs.slide_layouts[1]  # Title and content layout
        else:
            layout = prs.slide_layouts[0]  # Use first available layout
            
        slide = prs.slides.add_slide(layout)
        
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get('title', 'Untitled')
        
        # Add bullet points
        bullets = slide_data.get('bullets', [])
        if bullets:
            # Try to find content placeholder
            content_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1:  # Content placeholder
                    content_placeholder = shape
                    break
            
            if content_placeholder:
                tf = content_placeholder.text_frame
                tf.clear()
                for bullet in bullets:
                    p = tf.add_paragraph()
                    p.text = str(bullet)
                    p.level = 0
            else:
                # If no content placeholder, add text box
                left = Inches(1)
                top = Inches(2)
                width = Inches(8)
                height = Inches(4)
                text_box = slide.shapes.add_textbox(left, top, width, height)
                tf = text_box.text_frame
                for bullet in bullets:
                    p = tf.add_paragraph()
                    p.text = f"• {bullet}"
        
        # Speaker notes (optional)
        if 'notes' in slide_data and slide_data['notes']:
            slide.notes_slide.notes_text_frame.text = slide_data['notes']
    
    prs.save(output_path)
